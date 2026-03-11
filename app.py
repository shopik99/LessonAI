import streamlit as st
from google import genai
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from io import BytesIO
import requests
from pptx.enum.text import PP_ALIGN

# ====== Secrets.toml арқылы API кілтін алу ======
API_KEY = st.secrets["GEMINI_API_KEY"]
client = genai.Client(api_key=API_KEY)

st.set_page_config(page_title="Кәсіби презентация", page_icon="🎓")
st.title("🎓 Толық кәсіби презентация генераторы")

# Пайдаланушыдан пән, тақырып және слайд саны
subject = st.text_input("Пән атауы:")
topic = st.text_input("Сабақ тақырыбы:")
num_slides = st.number_input("Слайд саны (макс 10):", min_value=1, max_value=10, value=5, step=1)

if st.button("Презентация жасау") and subject and topic:
    prompt = f"""
    Маған '{subject}' пәні бойынша '{topic}' тақырыбында {num_slides} слайдтық презентация жасаңдар.
    Әр слайдқа:
    - Қазақша тақырып (қысқа және нақты)
    - Теориялық мәтін, кемінде 5 сөйлем
    - Тақырыпқа сай сурет URL
    Әр слайд бөлек жолға болсын, тек тақырып атауы және мәтін, "Слайд 1/2" деп жазбаңыз.
    """

    try:
        response = client.models.generate_content(
            model="gemini-2.5-flash",
            contents=prompt
        )
        slides_text = response.text
    except Exception as e:
        st.error(f"GPT сұранысында қате: {e}")
        slides_text = None

    if slides_text:
        prs = Presentation()
        background_colors = [(240, 248, 255), (255, 248, 220), (245, 255, 250), (255, 240, 245)]

        slide_index = 0
        # GPT жауапын әр жол слайд деп есептейміз
        slide_lines = [line.strip() for line in slides_text.split("\n") if line.strip()][:num_slides]

        for slide_line in slide_lines:
            # Слайд тақырыбы мен мәтінін бөліп алу
            if "|" in slide_line:  # мысалы GPT "Тақырып | Мәтін | ImageURL" форматында берсе
                parts = slide_line.split("|")
                slide_title = parts[0].strip()
                slide_content = parts[1].strip() if len(parts) > 1 else ""
                img_url = parts[2].strip() if len(parts) > 2 else None
            else:
                # Егер URL жоқ болса, тақырып пен мәтінді бөліп аламыз
                slide_title = f"Тақырып {slide_index+1}"
                slide_content = slide_line
                img_url = None

            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = slide_title

            # Фон түсі
            fill = slide.background.fill
            color = background_colors[slide_index % len(background_colors)]
            fill.solid()
            fill.fore_color.rgb = RGBColor(*color)

            # Мәтін paragraph-қа бөлу
            textbox = slide.placeholders[1]
            textbox.text = ""
            sentences = [s.strip() for s in slide_content.split(". ") if s.strip()]
            for s in sentences:
                p = textbox.text_frame.add_paragraph()
                p.text = s
                p.level = 0
                p.font.size = Pt(14)
                p.font.name = "Times New Roman"
                p.font.color.rgb = RGBColor(0,0,0)
                p.alignment = PP_ALIGN.LEFT

            # Сурет қосу
            if img_url:
                try:
                    img_data = requests.get(img_url).content
                    slide.shapes.add_picture(BytesIO(img_data), Inches(5.5), Inches(1.5), width=Inches(3))
                except:
                    pass

            slide_index += 1

        # PPTX файл байтқа жазу
        pptx_bytes = BytesIO()
        prs.save(pptx_bytes)
        pptx_bytes.seek(0)

        st.success("🎉 Кәсіби презентация дайын!")
        st.download_button(
            label="📥 PPTX жүктеу",
            data=pptx_bytes,
            file_name=f"{topic}_presentation.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
